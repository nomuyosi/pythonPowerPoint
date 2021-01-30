
#参考
#https://qiita.com/aykbzl/items/09b52fabb3af6b925fb9
#https://qiita.com/shimajiroxyz/items/4316608a01eb91543faa

#pip install python-pptx
#pip install Pillow

from pptx import Presentation
from PIL import Image
import os
import datetime

#スライドサイズ
#4:3 (default) 9144000x6858000
#16:9 12193200x6858000
SLIDE_WIDTH, SLIDE_HEIGHT = 12193200, 6858000
#スライド中心のX、Y座標（左上が原点）
IMG_CENTER_X, IMG_CENTER_Y = SLIDE_WIDTH/2, SLIDE_HEIGHT/2
#スライドのアスペクト比
SLIDE_ASPECT_RATIO = SLIDE_WIDTH / SLIDE_HEIGHT

#出力ファイル名
#OUTPUT_FILE_PATH = "test.pptx"
#画像の格納ディレクトリ
#IMG_DIR = "./data/"

#受け取ったプレゼンテーションオブジェクトにスライドを追加し、追加されたスライドオブジェクトを返す。
def add_slide(prs):
  #白紙スライドの追加(ID=6は白紙スライド)
  blank_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank_slide_layout)
  return slide

#画像をスライド中心に貼り付ける
def add_picture(slide, img_file):
  #画像サイズを取得してアスペクト比を得る
  im = Image.open(img_file)
  im_width, im_height = im.size
  aspect_ratio = im_width/im_height

  #スライドと画像のアスペクト比に応じて処理を分岐
  #画像のほうが横長だったら横めいっぱいに広げる
  if aspect_ratio > SLIDE_ASPECT_RATIO:
    img_display_width = SLIDE_WIDTH
    img_display_height = img_display_width / aspect_ratio
  else: #画像のほうが縦長だったら縦めいっぱいに広げる
    img_display_height = SLIDE_HEIGHT
    img_display_width = img_display_height * aspect_ratio
  #センタリングする場合の画像の左上座標を計算
  left = IMG_CENTER_X - img_display_width / 2
  top = IMG_CENTER_Y - img_display_height / 2

  #画像をスライドに追加
  if aspect_ratio > SLIDE_ASPECT_RATIO:
    slide.shapes.add_picture(img_file, left, top, width = img_display_width)
  else:
    slide.shapes.add_picture(img_file, left, top, height = img_display_height)

  return slide

def py_ppt_add_img(IMG_DIR):

  #スライドオブジェクトの定義
  prs = Presentation()
  #スライドサイズの指定
  prs.slide_width = SLIDE_WIDTH
  prs.slide_height = SLIDE_HEIGHT

  #img画像のファイル名を取得
  img_files = os.listdir(IMG_DIR)
  #pngで終了するファイル名のみ抽出。貼り付けたい画像の拡張子に応じて変える
  img_files = [name for name in img_files if name.endswith(".png")]
  #昇順にソート（この順番でスライドに貼り付けられる）
  img_files.sort()#昇順にsort

  for name in img_files:
    path = IMG_DIR + name
    slide = add_slide(prs)
    add_picture(slide, path)

  dt_now = datetime.datetime.now()
  print(dt_now.strftime('%Y年%m月%d日 %H:%M:%S'))
  d_today = dt_now.strftime('%Y%m%d_%H%M%S')
  #pptxファイルを出力する
  OUTPUT_FILE_PATH = IMG_DIR + d_today + "_output.pptx"
  print("output = ", OUTPUT_FILE_PATH)
  prs.save(OUTPUT_FILE_PATH)